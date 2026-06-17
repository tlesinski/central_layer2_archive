"""Generate PARTMGR Technical Overview PowerPoint presentation with diagrams."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── colour palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xE8, 0x6C, 0x00)
LIGHT_BG   = RGBColor(0xF0, 0xF2, 0xF5)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
GREY       = RGBColor(0x66, 0x66, 0x66)
CODE_BG    = RGBColor(0x28, 0x2C, 0x34)
CODE_FG    = RGBColor(0xAB, 0xB2, 0xBF)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
BLUE       = RGBColor(0x29, 0x80, 0xB9)
RED        = RGBColor(0xC0, 0x39, 0x2B)
YELLOW_BG  = RGBColor(0xF3, 0x9C, 0x12)
TEAL       = RGBColor(0x16, 0xA0, 0x85)
AGENT_C    = RGBColor(0x34, 0x98, 0xDB)
ARCH_C     = RGBColor(0x2E, 0xCC, 0x71)
REPL_C     = RGBColor(0x9B, 0x59, 0xB6)
CLIENT_C   = RGBColor(0xE6, 0x7E, 0x22)
TARGET_C   = RGBColor(0x1A, 0xBC, 0x9C)
LIGHT_GREY = RGBColor(0xEC, 0xF0, 0xF1)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── generic helpers ─────────────────────────────────────────────
def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_shape(slide, left, top, width, height, fill_color=None,
               line_color=None, line_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def _add_rounded_box(slide, left, top, width, height, fill_color, text,
                     font_size=14, font_color=WHITE, bold=True):
    shape = _add_shape(slide, left, top, width, height,
                       fill_color=fill_color,
                       shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    return shape

def _add_arrow_right(slide, left, top, width=Inches(0.6), height=Inches(0.3), color=GREY):
    shape = _add_shape(slide, left, top, width, height,
                       fill_color=color,
                       shape_type=MSO_SHAPE.RIGHT_ARROW)
    return shape

def _add_arrow_down(slide, left, top, width=Inches(0.3), height=Inches(0.5), color=GREY):
    shape = _add_shape(slide, left, top, width, height,
                       fill_color=color,
                       shape_type=MSO_SHAPE.DOWN_ARROW)
    return shape

def _add_textbox(slide, left, top, width, height, text="",
                 font_size=18, bold=False, color=DARK_GREY,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri",
                 anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    return txBox

def add_slide_number(slide, num, total=32):
    _add_textbox(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
                 f"{num} / {total}", font_size=10, color=GREY,
                 alignment=PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    _add_shape(slide, Inches(0), Inches(7.15), W, Inches(0.35),
               fill_color=NAVY)

def add_title_header(slide, title, subtitle=None):
    _add_shape(slide, Inches(0), Inches(0), W, Inches(1.15), fill_color=NAVY)
    _add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.65),
                 title, font_size=30, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(0.72), Inches(12), Inches(0.35),
                     subtitle, font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD))

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = _add_shape(slide, left, top, width, height, fill_color=CODE_BG)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), height - Inches(0.2),
                 code_text, font_size=font_size, color=CODE_FG,
                 font_name="Consolas")
    return shape

def add_bullets(slide, left, top, width, height, items, font_size=16,
                color=DARK_GREY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_legend_box(slide, left, top, width, height, items):
    """Add a color legend. items = [(color, label), ...]"""
    box = _add_shape(slide, left, top, width, height,
                     fill_color=RGBColor(0xFD, 0xFD, 0xFD),
                     line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
    y = top + Inches(0.08)
    for color, label in items:
        _add_shape(slide, left + Inches(0.1), y, Inches(0.2), Inches(0.2),
                   fill_color=color)
        _add_textbox(slide, left + Inches(0.38), y - Inches(0.02),
                     width - Inches(0.5), Inches(0.25),
                     label, font_size=10, color=GREY)
        y += Inches(0.28)

# ── diagram helpers ─────────────────────────────────────────────

# ── slide layout helpers ────────────────────────────────────────
def slide_section(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 title, font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                     subtitle, font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD),
                     alignment=PP_ALIGN.CENTER)
    add_bottom_bar(slide)
    return slide

def slide_content(title, bullets, subtitle=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_title_header(slide, title, subtitle)
    add_bullets(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.2), bullets)
    if note:
        _add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
                     note, font_size=11, color=GREY)
    add_bottom_bar(slide)
    add_slide_number(slide, len(prs.slides))
    return slide

def slide_with_code(title, bullets, code, subtitle=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_title_header(slide, title, subtitle)
    add_bullets(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.5),
                bullets, font_size=15)
    add_code_block(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
                   code, font_size=11)
    if note:
        _add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
                     note, font_size=11, color=GREY)
    add_bottom_bar(slide)
    add_slide_number(slide, len(prs.slides))
    return slide

def slide_full_code(title, code, subtitle=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_title_header(slide, title, subtitle)
    add_code_block(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.4),
                   code, font_size=10)
    if note:
        _add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
                     note, font_size=11, color=GREY)
    add_bottom_bar(slide)
    add_slide_number(slide, len(prs.slides))
    return slide

def slide_two_columns(title, left_bullets, right_bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_title_header(slide, title, subtitle)
    add_bullets(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.2),
                left_bullets, font_size=15)
    add_bullets(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
                right_bullets, font_size=15)
    add_bottom_bar(slide)
    add_slide_number(slide, len(prs.slides))
    return slide

def add_pipeline_flow(slide, left, top, width, steps, arrow_color=GREY):
    """Draw a horizontal pipeline with coloured boxes and arrows.
    steps = [(label, color), ...]
    """
    n = len(steps)
    box_w = width // n - Inches(0.35)
    if box_w < Inches(1.2):
        box_w = Inches(1.2)
    arrow_w = Inches(0.3)
    total_w = n * box_w + (n - 1) * arrow_w
    start_x = left + (width - total_w) // 2
    x = start_x

    for i, (label, color) in enumerate(steps):
        _add_rounded_box(slide, x, top, box_w, Inches(0.65), color,
                         label, font_size=13, font_color=WHITE, bold=True)
        x += box_w
        if i < n - 1:
            _add_arrow_right(slide, x, top + Inches(0.18), arrow_w, Inches(0.3),
                             color=arrow_color)
            x += arrow_w

def add_three_tier_diagram(slide, left, top, width, height):
    """3-tier architecture: 3 databases (TWP / TWARP / TWART) with components inside."""
    db_w = Inches(3.2)
    db_h = Inches(2.2)
    gap = Inches(0.35)
    total_w = 3 * db_w + 2 * gap
    start_x = left + (width - total_w) // 2
    mid_y = top

    db_color = RGBColor(0x2C, 0x3E, 0x50)
    inner_w = db_w - Inches(0.4)

    # ── DB 1: TWP ──
    x1 = start_x
    _add_rounded_box(slide, x1, mid_y, db_w, db_h, db_color,
                     "TWP (Source DB)", font_size=13, font_color=WHITE)

    _add_rounded_box(slide, x1 + Inches(0.2), mid_y + Inches(0.65), inner_w, Inches(0.6),
                     CLIENT_C, "CLIENT1 / CLIENT2", font_size=12)

    _add_rounded_box(slide, x1 + Inches(0.2), mid_y + Inches(1.4), inner_w, Inches(0.6),
                     AGENT_C, "PARTMGR_AGENT", font_size=12)

    # ── arrow 1 -> 2 ──
    _add_arrow_right(slide, x1 + db_w + Inches(0.04), mid_y + db_h//2 - Inches(0.1),
                     gap - Inches(0.08), Inches(0.25))

    # ── DB 2: TWARP ──
    x2 = x1 + db_w + gap
    _add_rounded_box(slide, x2, mid_y, db_w, db_h, db_color,
                     "TWARP (Archive DB)", font_size=13, font_color=WHITE)

    _add_rounded_box(slide, x2 + Inches(0.2), mid_y + Inches(0.8), inner_w, Inches(0.6),
                     ARCH_C, "PARTMGR_ARCHIVER", font_size=12)

    # ── arrow 2 -> 3 ──
    _add_arrow_right(slide, x2 + db_w + Inches(0.04), mid_y + db_h//2 - Inches(0.1),
                     gap - Inches(0.08), Inches(0.25))

    # ── DB 3: TWART ──
    x3 = x2 + db_w + gap
    _add_rounded_box(slide, x3, mid_y, db_w, db_h, db_color,
                     "TWART (Replica DB)", font_size=13, font_color=WHITE)

    _add_rounded_box(slide, x3 + Inches(0.2), mid_y + Inches(0.8), inner_w, Inches(0.6),
                     REPL_C, "PARTMGR_REPLICA", font_size=12)

def add_partition_lifecycle_diagram(slide, left, top, width, height):
    """Partition lifecycle: TWP(small window) → TWARP(full 2yr history) → TWART(6mo rolling)."""
    swimlane_h = Inches(1.1)
    gap = Inches(0.3)
    total_h = 3 * swimlane_h + 2 * gap
    start_y = top + (height - total_h) // 2

    label_w = Inches(1.8)
    tl_left = left + label_w + Inches(0.15)
    tl_w = width - (label_w + Inches(0.15))

    months = 24
    seg_w = tl_w // months

    # key transition points (in months)
    ARCHIVE_AT  = 1     # archive happens right after month 1
    TRUNC_AT    = 2     # truncate ~15 days after archive = month 2
    REPL_AT     = 18    # replicate to replica (rolling window)
    PURGE_BEFORE = 6    # only keep 6 months in replica

    online_color = RGBColor(0x5D, 0xAD, 0xEC)  # soft blue for "active"
    arch_color   = RGBColor(0x2E, 0xCC, 0x71)  # green for archived
    repl_color   = RGBColor(0x9B, 0x59, 0xB6)  # purple for replicated
    greyed       = RGBColor(0xCC, 0xCC, 0xCC)  # truncated/purged

    def m_x(m):
        return int(tl_left + seg_w * m)

    def draw_swimlane_bg(y, lbl, lbl_color):
        _add_shape(slide, left, y, width, swimlane_h,
                   fill_color=RGBColor(0xF8, 0xF9, 0xFA),
                   line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
        _add_rounded_box(slide, left + Inches(0.05), y + Inches(0.1),
                         label_w - Inches(0.1), swimlane_h - Inches(0.2),
                         lbl_color, lbl, font_size=11, font_color=WHITE)

    # ── TWP row ──
    y0 = start_y
    draw_swimlane_bg(y0, "TWP (CLIENT / AGENT)", CLIENT_C)
    # active blocks
    for m in range(0, ARCHIVE_AT + 1):
        bw = max(seg_w - 1, 1)
        _add_shape(slide, m_x(m), y0 + Inches(0.4), bw, Inches(0.5),
                   fill_color=online_color)
    # truncated (grey) after TRUNC_AT
    for m in range(TRUNC_AT, months):
        bw = max(seg_w - 1, 1)
        _add_shape(slide, m_x(m), y0 + Inches(0.4), bw, Inches(0.5),
                   fill_color=greyed)

    # TRUNCATE label on TWP
    trunc_x = m_x(TRUNC_AT)
    _add_textbox(slide, trunc_x - Inches(0.1), y0 - Inches(0.05), Inches(1.2), Inches(0.35),
                 "TRUNCATE\n(~15 days)", font_size=9, bold=True, color=RED,
                 alignment=PP_ALIGN.CENTER)
    # red X icon
    _add_textbox(slide, trunc_x + Inches(1.0), y0 + Inches(0.15), Inches(0.5), Inches(0.5),
                 "\u2718", font_size=18, bold=True, color=RED, alignment=PP_ALIGN.CENTER)
    # short retention label
    _add_textbox(slide, left + Inches(1.9), y0 + Inches(0.75), Inches(3.0), Inches(0.25),
                 "Short source retention (\u226415 days after archive)",
                 font_size=8, color=GREY, alignment=PP_ALIGN.CENTER)

    # ── TWP→TWARP: ARCHIVE arrow ──
    arch_x = m_x(ARCHIVE_AT) + seg_w // 2 - Inches(0.12)
    _add_down_arrow(slide, arch_x, y0 + swimlane_h - Inches(0.05),
                    Inches(0.25), gap + Inches(0.1), ARCH_C)
    _add_textbox(slide, arch_x - Inches(0.9), y0 + swimlane_h - Inches(0.15),
                 Inches(1.8), Inches(0.4), "ARCHIVE (EXCHANGE)",
                 font_size=9, bold=True, color=ARCH_C, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, arch_x - Inches(0.6), y0 + swimlane_h + Inches(0.15),
                 Inches(1.2), Inches(0.2), "DISCOVER\u2192ARCHIVE\u2192QUALITY",
                 font_size=7, color=GREY, alignment=PP_ALIGN.CENTER)

    # ── TWARP row ──
    y1 = start_y + swimlane_h + gap
    draw_swimlane_bg(y1, "TWARP (ARCHIVER)", ARCH_C)
    # full history - all months archived
    for m in range(0, months):
        bw = max(seg_w - 1, 1)
        _add_shape(slide, m_x(m), y1 + Inches(0.4), bw, Inches(0.5),
                   fill_color=arch_color)
    # full history annotation
    _add_textbox(slide, left + Inches(2.0), y1 + Inches(0.75), Inches(4.0), Inches(0.25),
                 "Full archive history (e.g. 2 years  \u2014  never truncated)",
                 font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── TWARP→TWART: REPLICATE arrow and rolling window indicator ──
    repl_x = m_x(REPL_AT) + seg_w // 2 - Inches(0.12)
    _add_down_arrow(slide, repl_x, y1 + swimlane_h - Inches(0.05),
                    Inches(0.25), gap + Inches(0.1), REPL_C)
    _add_textbox(slide, repl_x - Inches(1.1), y1 + swimlane_h - Inches(0.15),
                 Inches(2.0), Inches(0.4), "REPLICATE\n(rolling window feed)",
                 font_size=9, bold=True, color=REPL_C, alignment=PP_ALIGN.CENTER)

    # rolling window bracket / annotation on TWARP side
    _add_textbox(slide, repl_x - Inches(1.6), y1 + Inches(0.1), Inches(1.5), Inches(0.3),
                 "Qualified\npartitions only", font_size=7, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    # ── TWART row ──
    y2 = start_y + 2 * (swimlane_h + gap)
    draw_swimlane_bg(y2, "TWART (REPLICA)", REPL_C)
    # purged (grey) for older months
    for m in range(0, months - PURGE_BEFORE):
        bw = max(seg_w - 1, 1)
        _add_shape(slide, m_x(m), y2 + Inches(0.4), bw, Inches(0.5),
                   fill_color=greyed)
    # active rolling window (last 6 months)
    for m in range(months - PURGE_BEFORE, months):
        bw = max(seg_w - 1, 1)
        _add_shape(slide, m_x(m), y2 + Inches(0.4), bw, Inches(0.5),
                   fill_color=repl_color)

    # PURGE label on TWART
    purge_x = m_x(months - PURGE_BEFORE)
    _add_textbox(slide, purge_x - Inches(0.2), y2 - Inches(0.05), Inches(1.2), Inches(0.35),
                 "PURGE", font_size=9, bold=True, color=RED,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, purge_x + Inches(1.0), y2 + Inches(0.15), Inches(0.5), Inches(0.5),
                 "\u2718", font_size=18, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    # rolling window annotation
    _add_textbox(slide, m_x(months - PURGE_BEFORE) - Inches(0.5), y2 + Inches(0.75),
                 Inches(3.5), Inches(0.25),
                 "Rolling 6-month window  \u2192  oldest partitions purged",
                 font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)

    # time axis labels
    halves = [(0, "2024 H1"), (6, "2024 H2"), (12, "2025 H1"), (18, "2025 H2")]
    for m_offset, lbl in halves:
        _add_textbox(slide, m_x(m_offset) - Inches(0.1),
                     start_y + 3 * (swimlane_h + gap) - Inches(0.05),
                     Inches(1.5), Inches(0.25),
                     lbl, font_size=9, bold=True, color=DARK_GREY,
                     alignment=PP_ALIGN.CENTER)

def _add_down_arrow(slide, left, top, width, height, color):
    shape = _add_shape(slide, left, top, width, height,
                       fill_color=color,
                       shape_type=MSO_SHAPE.DOWN_ARROW)
    return shape

def add_archiver_flow_diagram(slide, left, top, width, height):
    """Detailed ARCHIVER data flow with CLIENT -> AGENT -> ARCHIVER target"""
    # Source side
    box_w = Inches(2.0)
    box_h = Inches(0.85)

    _add_rounded_box(slide, left, top, box_w, box_h,
                     CLIENT_C, "CLIENT\nPartitioned Table", font_size=13)
    _add_textbox(slide, left, top + box_h + Inches(0.05), box_w, Inches(0.3),
                 "Source data", font_size=9, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    _add_arrow_right(slide, left + box_w + Inches(0.1), top + Inches(0.25),
                     Inches(0.4), Inches(0.3))

    x2 = left + box_w + Inches(0.6)
    _add_rounded_box(slide, x2, top, box_w, box_h,
                     AGENT_C, "AGENT\nVW_AGENT_PARTITION_INFO\nfn_get_row_count", font_size=12)
    _add_textbox(slide, x2, top + box_h + Inches(0.05), box_w, Inches(0.3),
                 "Metadata + row counts", font_size=9, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    _add_arrow_right(slide, x2 + box_w + Inches(0.1), top + Inches(0.25),
                     Inches(0.4), Inches(0.3))

    x3 = x2 + box_w + Inches(0.6)
    _add_rounded_box(slide, x3, top, box_w, box_h,
                     ARCH_C, "ARCHIVER\nTBL_ARCHIVER_TABLES\nTBL_ARCHIVER_PARTITIONS",
                     font_size=12)
    _add_textbox(slide, x3, top + box_h + Inches(0.05), box_w, Inches(0.3),
                 "Configuration + status", font_size=9, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    _add_arrow_right(slide, x3 + box_w + Inches(0.1), top + Inches(0.25),
                     Inches(0.4), Inches(0.3))

    x4 = x3 + box_w + Inches(0.6)
    _add_rounded_box(slide, x4, top, box_w, box_h,
                     TARGET_C, "ARCHIVER\nTarget Table\n(partitioned)", font_size=12)
    _add_textbox(slide, x4, top + box_h + Inches(0.05), box_w, Inches(0.3),
                 "Archived data", font_size=9, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    # DB link label
    _add_textbox(slide, x2 + Inches(0.2), top - Inches(0.3), Inches(1.5), Inches(0.25),
                 "DB link", font_size=8, color=GREY, alignment=PP_ALIGN.CENTER)

def add_pipeline_steps_diagram(slide, left, top, width):
    """Vertical pipeline with boxes and down arrows."""
    steps = [
        ("1. DISCOVER", "Read source partition metadata\nCreate missing target partitions", AGENT_C),
        ("2. ARCHIVE", "Create staging table\nLoad data via INSERT APPEND\nEXCHANGE PARTITION", ARCH_C),
        ("3. QUALITY", "Compare SOURCE_ROW_COUNT\nvs TARGET_ROW_COUNT", ORANGE),
        ("4. TRUNCATE\n(preview)", "Apply DAYS_ONLINE + PRESERVE_RULE\nCall AGENT cleanup\npreview-first", RED),
    ]
    box_w = Inches(2.5)
    box_h = Inches(1.2)
    gap = Inches(0.3)
    total_h = len(steps) * (box_h + gap) - gap
    start_y = top

    x_desc = left + box_w + Inches(0.4)
    desc_w = Inches(5.5)

    for i, (title, desc, color) in enumerate(steps):
        y = start_y + i * (box_h + gap)
        _add_rounded_box(slide, left, y, box_w, box_h, color,
                         title, font_size=14, font_color=WHITE)
        _add_textbox(slide, x_desc, y + Inches(0.15), desc_w, box_h - Inches(0.3),
                     desc, font_size=14, color=DARK_GREY)
        if i < len(steps) - 1:
            _add_arrow_down(slide, left + box_w//2 - Inches(0.15),
                            y + box_h, Inches(0.3), gap,
                            color=GREY)

def add_splited_topology_diagram(slide, left, top, width, height):
        _add_textbox(slide, left, top, width, Inches(0.4),
                     "SPLIT Model", font_size=16, bold=True, color=NAVY,
                     alignment=PP_ALIGN.CENTER)
        # three separate boxes
        y = top + Inches(0.6)
        # DB1: TWP
        _add_rounded_box(slide, left + Inches(0.1), y, Inches(3.0), Inches(1.8),
                         RGBColor(0x1A, 0x52, 0x7A),
                         "TWP (Source DB)", font_size=14)
        _add_rounded_box(slide, left + Inches(0.3), y + Inches(0.5),
                         Inches(2.6), Inches(0.55),
                         CLIENT_C, "CLIENT1 / CLIENT2", font_size=11)
        _add_rounded_box(slide, left + Inches(0.3), y + Inches(1.15),
                         Inches(2.6), Inches(0.5),
                         AGENT_C, "PARTMGR_AGENT", font_size=11)
        _add_textbox(slide, left + Inches(0.1), y + Inches(1.85),
                     Inches(3.0), Inches(0.3),
                     "SYS: SOURCE_SYS_CONNECT", font_size=8, color=GREY,
                     alignment=PP_ALIGN.CENTER)
        # arrow
        _add_arrow_right(slide, left + Inches(3.2), y + Inches(0.7),
                         Inches(0.3), Inches(0.3))
        # DB2: TWARP
        _add_rounded_box(slide, left + Inches(3.6), y, Inches(3.0), Inches(1.8),
                         RGBColor(0x1A, 0x6B, 0x52),
                         "TWARP (Archive DB)", font_size=14)
        _add_rounded_box(slide, left + Inches(3.8), y + Inches(0.8),
                         Inches(2.6), Inches(0.55),
                         ARCH_C, "PARTMGR_ARCHIVER", font_size=11)
        _add_textbox(slide, left + Inches(3.6), y + Inches(1.85),
                     Inches(3.0), Inches(0.3),
                     "SYS: ARCHIVER_SYS_CONNECT", font_size=8, color=GREY,
                     alignment=PP_ALIGN.CENTER)
        # arrow
        _add_arrow_right(slide, left + Inches(6.7), y + Inches(0.7),
                         Inches(0.3), Inches(0.3))
        # DB3: TWART
        _add_rounded_box(slide, left + Inches(7.1), y, Inches(3.0), Inches(1.8),
                         RGBColor(0x6C, 0x34, 0x7A),
                         "TWART (Replica DB)", font_size=14)
        _add_rounded_box(slide, left + Inches(7.3), y + Inches(0.8),
                         Inches(2.6), Inches(0.55),
                         REPL_C, "PARTMGR_REPLICA", font_size=11)
        _add_textbox(slide, left + Inches(7.1), y + Inches(1.85),
                     Inches(3.0), Inches(0.3),
                     "SYS: REPLICA_SYS_CONNECT", font_size=8, color=GREY,
                     alignment=PP_ALIGN.CENTER)

        # link labels
        _add_textbox(slide, left + Inches(3.2), y + Inches(0.3),
                     Inches(0.5), Inches(0.3),
                     "DB\nlink", font_size=7, color=GREY, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, left + Inches(6.7), y + Inches(0.3),
                     Inches(0.5), Inches(0.3),
                     "DB\nlink", font_size=7, color=GREY, alignment=PP_ALIGN.CENTER)


def add_replica_flow_diagram(slide, left, top, width, height):
    """REPLICA flow: ARCHIVER -> DB link -> REPLICA -> target"""
    box_w = Inches(2.0)
    box_h = Inches(0.85)
    mid_y = top + height // 2

    _add_rounded_box(slide, left, mid_y - box_h//2, box_w, box_h,
                     ARCH_C, "ARCHIVER\nQualified Partitions", font_size=12)
    _add_textbox(slide, left, mid_y + box_h//2 + Inches(0.05), box_w, Inches(0.25),
                 "archive_status + quality_status = Y", font_size=8, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    _add_arrow_right(slide, left + box_w + Inches(0.1), mid_y - Inches(0.15),
                     Inches(0.4), Inches(0.3))

    x2 = left + box_w + Inches(0.6)
    _add_rounded_box(slide, x2, mid_y - box_h//2, box_w, box_h,
                     REPL_C, "REPLICA\nMetadata + Staging", font_size=12)
    _add_textbox(slide, x2, mid_y + box_h//2 + Inches(0.05), box_w, Inches(0.25),
                 "REPLICA_ARCHIVER_PARTITIONS_SRC", font_size=8, color=GREY,
                 alignment=PP_ALIGN.CENTER)

    _add_arrow_right(slide, x2 + box_w + Inches(0.1), mid_y - Inches(0.15),
                     Inches(0.4), Inches(0.3))

    x3 = x2 + box_w + Inches(0.6)
    _add_rounded_box(slide, x3, mid_y - box_h//2, box_w, box_h,
                     TARGET_C, "REPLICA\nTarget Table\n(partitioned)", font_size=12)
    _add_textbox(slide, x3, mid_y + box_h//2 + Inches(0.05), box_w, Inches(0.25),
                 "Replicated data", font_size=8, color=GREY,
                 alignment=PP_ALIGN.CENTER)


# ── Build slides ────────────────────────────────────────────────

# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, NAVY)
_add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.0),
             "PARTMGR", font_size=54, bold=True, color=ORANGE,
             alignment=PP_ALIGN.CENTER)
_add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.8),
             "Oracle Archive Control Plane", font_size=36, bold=False, color=WHITE,
             alignment=PP_ALIGN.CENTER)
_add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "Central Layer 2 archive management for partitioned Oracle estates",
             font_size=18, color=RGBColor(0xBB, 0xCC, 0xDD),
             alignment=PP_ALIGN.CENTER)
_add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "Tomasz Lesi\u0144ski  |  2026-06", font_size=16, color=GREY,
             alignment=PP_ALIGN.CENTER)
add_bottom_bar(slide)
add_slide_number(slide, 1)

# ── SLIDE 2: Problem ──
slide_content(
    "The Problem",
    [
        "Large Oracle systems have hundreds of partitioned tables across many schemas",
        "Each source uses its own, often ad-hoc, archive logic",
        "No central visibility into what has been archived and what hasn't",
        "No quality guarantees -- source and target row counts frequently diverge",
        "Source cleanup is risky and irreversible",
        "Downstream systems (reporting, audit) need bounded copies of archived data"
    ],
    subtitle="Challenges in managing partitioned data archiving",
    note="PARTMGR solves these problems by centralizing archive orchestration outside source applications."
)

# ── SLIDE 3: Components Overview ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, WHITE)
add_title_header(slide, "System Components",
                 "Three layers: AGENT \u00b7 ARCHIVER \u00b7 REPLICA")
add_three_tier_diagram(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.0))

# Legend and descriptions
add_legend_box(slide, Inches(0.5), Inches(5.0), Inches(3.0), Inches(1.8), [
    (CLIENT_C, "Source business schemas"),
    (AGENT_C, "Layer 1: source-side helper"),
    (ARCH_C, "Layer 2: central control plane"),
    (REPL_C, "Layer 3: downstream replica"),
])
add_bullets(slide, Inches(4.0), Inches(5.0), Inches(8.5), Inches(1.8), [
    "AGENT: partition metadata, row counts, cleanup (no policy)",
    "ARCHIVER: orchestration, configuration, quality, truncate preview",
    "REPLICA: optional bounded downstream copy",
], font_size=14)
add_bottom_bar(slide)
add_slide_number(slide, len(prs.slides))

# ── SLIDE 4: Partition Lifecycle ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, WHITE)
add_title_header(slide, "Partition Lifecycle Across Layers",
                 "Source \u2192 Archive \u2192 Replica with truncate and purge timelines")
add_partition_lifecycle_diagram(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.6))
add_bottom_bar(slide)
add_slide_number(slide, len(prs.slides))

# ── SLIDE 5: Core Goal ──
slide_content(
    "Core Goal & Design Principles",
    [
        "Centralize archive logic outside source applications",
        "Every archive unit is explicitly tracked in metadata",
        "Component communication through real Oracle database links",
        "Source vs target row counts are compared before any cleanup",
        "Destructive operations default to preview mode (preview-first)",
        "Minimum viable deployment: AGENT + ARCHIVER (REPLICA is optional)"
    ],
    subtitle="Key design tenets",
    note="Source: README.md, docs/architecture.md"
)

# ── SLIDE 6: Source Identity ──
slide_with_code(
    "Source Identity",
    [
        "Source identity key (composite primary key):",
        "  SOURCE_DB_LINK + SOURCE_OWNER + SOURCE_TABLE_NAME",
        "",
        "SOURCE_DB_LINK must be a real database link",
        "Values LOCAL, NONE, NULL are forbidden",
        "",
        "No surrogate keys (no archive IDs)",
        "No ARCHIVE_METHOD column",
        "No source registry table",
        "",
        "One row = one source mapping"
    ],
    "-- TBL_ARCHIVER_TABLES configuration\n"
    "INSERT INTO TBL_ARCHIVER_TABLES (\n"
    "  SOURCE_DB_LINK,\n"
    "  SOURCE_OWNER,\n"
    "  SOURCE_TABLE_NAME,\n"
    "  TARGET_OWNER,\n"
    "  TARGET_TABLE_NAME,\n"
    "  LAST_BUSINESS_DATE,\n"
    "  DAYS_ONLINE,\n"
    "  ENABLED_FLAG\n"
    ") VALUES (\n"
    "  'AGENT_LINK',\n"
    "  'CLIENT1',\n"
    "  'ORDERS_ARCH_SRC',\n"
    "  'PARTMGR_ARCHIVER',\n"
    "  'TBL_ARCHIVER_CLIENT1_ARCH',\n"
    "  DATE '2026-06-01',\n"
    "  30,\n"
    "  'Y'\n"
    ");",
    note="Source: AGENTS.md -- Central Model Rules"
)

# ── SLIDE 7: Partition Detection ──
slide_content(
    "Partition Detection -- HIGH_VALUE",
    [
        "Identity is based on partition/subpartition HIGH_VALUE, not the physical name",
        "Physical names (SYS_P...) are unstable -- they change after table rebuilds",
        "HIGH_VALUE is a DDL fragment, e.g. TO_DATE('2024-02-01','YYYY-MM-DD')",
        "",
        "AGENT queries ALL_TAB_PARTITIONS / ALL_TAB_SUBPARTITIONS",
        "HIGH_VALUE is converted from LONG to VARCHAR2 via:",
        "  \u2022 DBMS_LOB.SUBSTR(TO_CLOB(high_value)) -- in the package body",
        "  \u2022 EXTRACTVALUE(DBMS_XMLGEN.GETXMLTYPE(...)) -- in VW_AGENT_PARTITION_INFO",
        "",
        "ARCHIVER detects new partitions by comparing against TBL_ARCHIVER_PARTITIONS"
    ],
    subtitle="Identification by HIGH_VALUE, not by physical name"
)

# ── SLIDE 8: Supported Partition Types ──
slide_content(
    "Supported Partitioning Schemes",
    [
        "RANGE -- range partitioning on a DATE column",
        "RANGE-LIST -- range partitioning with LIST subpartitions (e.g. by STATUS_CODE)",
        "INTERVAL -- automatic partition creation on data insert",
        "INTERVAL-LIST -- interval with LIST subpartitions",
        "",
        "Requirement: the partition key must be a DATE column",
        "Partition boundaries must be time-based",
        "MAXVALUE is supported (automatically excluded from archiving)",
        "",
        "Type is determined by ARCHIVE_UNIT_TYPE:",
        "  \u2022 PARTITION -- for simple partitions",
        "  \u2022 SUBPARTITION -- for subpartition-level units"
    ],
    subtitle="RANGE \u00b7 RANGE-LIST \u00b7 INTERVAL \u00b7 INTERVAL-LIST"
)

# ── SLIDE 9: Example Table ──
slide_content(
    "Example Source Table",
    [
        "Table: CLIENT1.ORDERS_SUBPART_SRC",
        "Partitioning: RANGE(ORDER_DATE) with LIST(STATUS_CODE) subpartitions",
        "",
        "Source partitions:",
        "  P_ERROR (< 1800-01-01), P202401, P202402, P202403, P202404",
        "  P202512, P202601, P202602, P202603, P202604, P202605",
        "  PMAX (MAXVALUE)",
        "",
        "Subpartitions: SP_OPEN (NEW, PAID), SP_DONE (SHIPPED, CLOSED)",
        "",
        "Data flow:",
        "  CLIENT1 \u2192 AGENT_LINK \u2192 PARTMGR_ARCHIVER \u2192 ARCHIVER_LINK \u2192 PARTMGR_REPLICA"
    ],
    subtitle="CLIENT1.ORDERS_SUBPART_SRC -- RANGE-LIST"
)

# ── SLIDE 10: DDL Example ──
slide_full_code(
    "Source Table DDL -- Example",
    "CREATE TABLE CLIENT1.ORDERS_SUBPART_SRC (\n"
    "  ORDER_ID     NUMBER        NOT NULL,\n"
    "  CUSTOMER_ID  NUMBER        NOT NULL,\n"
    "  ORDER_DATE   DATE          NOT NULL,\n"
    "  STATUS_CODE  VARCHAR2(20)  NOT NULL,\n"
    "  REGION_CODE  VARCHAR2(10)  NOT NULL,\n"
    "  AMOUNT       NUMBER(12,2)  NOT NULL,\n"
    "  CREATED_AT   TIMESTAMP     NOT NULL\n"
    ")\n"
    "PARTITION BY RANGE (ORDER_DATE)\n"
    "SUBPARTITION BY LIST (STATUS_CODE)\n"
    "SUBPARTITION TEMPLATE (\n"
    "  SUBPARTITION SP_OPEN VALUES ('NEW', 'PAID'),\n"
    "  SUBPARTITION SP_DONE VALUES ('SHIPPED', 'CLOSED')\n"
    ")\n"
    "(\n"
    "  PARTITION P_ERROR VALUES LESS THAN (DATE '1800-01-01'),\n"
    "  PARTITION P202401 VALUES LESS THAN (DATE '2024-02-01'),\n"
    "  PARTITION P202402 VALUES LESS THAN (DATE '2024-03-01'),\n"
    "  PARTITION P202403 VALUES LESS THAN (DATE '2024-04-01'),\n"
    "  PARTITION P202404 VALUES LESS THAN (DATE '2024-05-01'),\n"
    "  PARTITION P202512 VALUES LESS THAN (DATE '2026-01-01'),\n"
    "  PARTITION PMAX  VALUES LESS THAN (MAXVALUE)\n"
    ");",
    subtitle="RANGE-LIST with STATUS_CODE subpartitioning",
    note="Partition key: ORDER_DATE (DATE). Subpartitions: LIST(STATUS_CODE)."
)

# ── SLIDE 11: ARCHIVER Data Flow (with diagram) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, WHITE)
add_title_header(slide, "ARCHIVER Data Flow",
                 "From source to archive through real DB links")
add_archiver_flow_diagram(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.5))
add_bullets(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(2.0), [
    "Metadata: VW_AGENT_PARTITION_INFO@AGENT_LINK \u2192 VW_ARCHIVER_SOURCE_PARTITIONS \u2192 VW_ARCHIVER_DISCOVERY_PARTITIONS \u2192 TBL_ARCHIVER_PARTITIONS",
    "Data: CLIENT partition \u2192 staging table (INSERT APPEND) \u2192 EXCHANGE PARTITION \u2192 ARCHIVER target",
    "Row counts: PKG_AGENT_ARCHIVE.fn_get_row_count@AGENT_LINK compared with TARGET_ROW_COUNT",
], font_size=14)
add_bottom_bar(slide)
add_slide_number(slide, len(prs.slides))

# ── SLIDE 12: ARCHIVER Metadata ──
slide_content(
    "ARCHIVER Metadata Tables",
    [
        "TBL_ARCHIVER_TABLES -- source-to-target configuration:",
        "  SOURCE_DB_LINK, SOURCE_OWNER, SOURCE_TABLE_NAME",
        "  TARGET_OWNER, TARGET_TABLE_NAME",
        "  LAST_BUSINESS_DATE -- archive upper boundary",
        "  DAYS_ONLINE -- retention window before source cleanup",
        "  PRESERVE_RULE -- SQL returning dates that must stay online",
        "  ENABLED_FLAG -- enables/disables the source",
        "",
        "TBL_ARCHIVER_PARTITIONS -- per-unit status:",
        "  ARCHIVE_STATUS, QUALITY_STATUS, TRUNCATE_STATUS",
        "  SOURCE_ROW_COUNT, TARGET_ROW_COUNT",
        "",
        "TBL_ARCHIVER_RUNS -- process run tracking",
        "TBL_ARCHIVER_PROCESS_LOG -- detailed process logs"
    ],
    subtitle="Four core metadata tables",
    note="PK: (SOURCE_DB_LINK, SOURCE_OWNER, SOURCE_TABLE_NAME, PARTITION_HIGH_VALUE, SUBPARTITION_HIGH_VALUE)"
)

# ── SLIDE 13: ARCHIVER Target DDL ──
slide_full_code(
    "ARCHIVER Target Table DDL",
    "CREATE TABLE PARTMGR_ARCHIVER.TBL_ARCHIVER_CLIENT1_SUBPART (\n"
    "  ORDER_ID     NUMBER        NOT NULL,\n"
    "  CUSTOMER_ID  NUMBER        NOT NULL,\n"
    "  ORDER_DATE   DATE          NOT NULL,\n"
    "  STATUS_CODE  VARCHAR2(20)  NOT NULL,\n"
    "  REGION_CODE  VARCHAR2(10)  NOT NULL,\n"
    "  AMOUNT       NUMBER(12,2)  NOT NULL,\n"
    "  CREATED_AT   TIMESTAMP     NOT NULL\n"
    ")\n"
    "PARTITION BY RANGE (ORDER_DATE)\n"
    "SUBPARTITION BY LIST (STATUS_CODE)\n"
    "SUBPARTITION TEMPLATE (\n"
    "  SUBPARTITION SP_OPEN VALUES ('NEW', 'PAID'),\n"
    "  SUBPARTITION SP_DONE VALUES ('SHIPPED', 'CLOSED')\n"
    ")\n"
    "(\n"
    "  PARTITION P_ERROR VALUES LESS THAN (DATE '1800-01-01')\n"
    ");",
    subtitle="Target table structure mirrors the source",
    note="ARCHIVER adds data partitions dynamically during DISCOVER. P_ERROR is the initial partition."
)

# ── SLIDE 14: Config Example ──
slide_with_code(
    "Configuration Example",
    [
        "Source: CLIENT1.ORDERS_SUBPART_SRC",
        "Target: PARTMGR_ARCHIVER.TBL_ARCHIVER_CLIENT1_SUBPART",
        "DB link: AGENT_LINK",
        "",
        "LAST_BUSINESS_DATE = DAT.fn_eod",
        "  \u2192 only partitions before this date are archived",
        "",
        "DAYS_ONLINE = 30",
        "  \u2192 source retention window before cleanup eligibility",
        "",
        "PRESERVE_RULE:",
        "  SELECT DAT.fn_boy FROM dual",
        "  UNION ALL",
        "  SELECT DAT.fn_eoy FROM dual",
        "  \u2192 year start and year end stay in the source"
    ],
    "-- TBL_ARCHIVER_TABLES row\n"
    "SOURCE_DB_LINK     = AGENT_LINK\n"
    "SOURCE_OWNER       = CLIENT1\n"
    "SOURCE_TABLE_NAME  = ORDERS_SUBPART_SRC\n"
    "TARGET_OWNER       = PARTMGR_ARCHIVER\n"
    "TARGET_TABLE_NAME  = TBL_ARCHIVER_CLIENT1_SUBPART\n"
    "LAST_BUSINESS_DATE = DAT.fn_eod\n"
    "DAYS_ONLINE        = 30\n"
    "PRESERVE_RULE      = SELECT DATE '2026-01-01'\n"
    "                       FROM dual\n"
    "                     UNION ALL\n"
    "                     SELECT DATE '2026-12-31'\n"
    "                       FROM dual\n"
    "ENABLED_FLAG       = Y",
    note="DAT.fn_eod returns DATE '2026-06-01' in the seed. In production, this is the actual business date."
)

# ── SLIDE 15: DISCOVER ──
slide_content(
    "DISCOVER -- Detecting New Partitions",
    [
        "1. Read VW_ARCHIVER_SOURCE_PARTITIONS (metadata via DB link)",
        "2. Compare against existing TBL_ARCHIVER_PARTITIONS rows",
        "3. For each new partition:",
        "   \u2022 ALTER TABLE ... ADD PARTITION ... -- create target partition",
        "   \u2022 INSERT INTO TBL_ARCHIVER_PARTITIONS -- add metadata row",
        "",
        "MAXVALUE and P_ERROR partitions are automatically excluded",
        "Supports both PARTITION and SUBPARTITION (ARCHIVE_UNIT_TYPE)",
        "",
        "Preview mode: p_execute = 'N' shows generated DDL",
        "Execute mode: p_execute = 'Y' runs DDL and inserts metadata"
    ],
    subtitle="Synchronizing source structure with archive targets",
    note="PKG_ARCHIVER_DISCOVERY.prc_discover(p_execute, p_target_owner, p_target_table_name)"
)

# ── SLIDE 16: ARCHIVE ──
slide_content(
    "ARCHIVE -- Loading Data",
    [
        "Selection: units with ARCHIVE_STATUS = 'N' and HIGH_VALUE < LAST_BUSINESS_DATE",
        "",
        "For each unit:",
        "  1. Create staging table (FOR EXCHANGE WITH TABLE)",
        "  2. Load data via INSERT APPEND from source DB link (with date filter)",
        "  3. Build local indexes on the staging table",
        "  4. Execute EXCHANGE PARTITION/SUBPARTITION with INCLUDING INDEXES",
        "  5. DROP staging table PURGE",
        "  6. Set ARCHIVE_STATUS = 'Y' and record TARGET_ROW_COUNT",
        "",
        "Parallel DML support (PARALLEL_DEGREE from configuration)",
        "",
        "Automatic cleanup of orphan staging tables older than N days"
    ],
    subtitle="Staging + EXCHANGE PARTITION for efficient loading",
    note="PKG_ARCHIVER_IMPORT.prc_import(p_execute, p_target_owner, p_target_table_name)"
)

# ── SLIDE 17: QUALITY ──
slide_content(
    "QUALITY -- Row Count Validation",
    [
        "Selection: units with ARCHIVE_STATUS = 'Y' and QUALITY_STATUS = 'N'",
        "",
        "For each unit:",
        "  1. Call PKG_AGENT_ARCHIVE.FN_GET_ROW_COUNT@DB_LINK",
        "     \u2192 source row count (queried via AGENT DB link)",
        "  2. Read TARGET_ROW_COUNT from TBL_ARCHIVER_PARTITIONS",
        "     \u2192 archive target row count",
        "  3. Compare:",
        "     \u2022 Match      \u2192 QUALITY_STATUS = 'Y'",
        "     \u2022 Mismatch   \u2192 QUALITY_STATUS = 'N' + ERROR_MESSAGE",
        "",
        "No quality pass = no truncate eligibility",
        "Generates summary: count of OK units vs errors"
    ],
    subtitle="Source-target row count comparison",
    note="PKG_ARCHIVER_QUALITY.prc_quality(p_execute, p_target_owner, p_target_table_name)"
)

# ── SLIDE 18: TRUNCATE ──
slide_content(
    "TRUNCATE -- Source Cleanup (Preview-First)",
    [
        "Selection: ARCHIVE_STATUS='Y' + QUALITY_STATUS='Y' + TRUNCATE_STATUS='N'",
        "Additional filters:",
        "  \u2022 HIGH_DATE <= LAST_BUSINESS_DATE - DAYS_ONLINE (retention window)",
        "  \u2022 PRESERVE_RULE -- SQL returning dates to keep online",
        "",
        "For each unit:",
        "  \u2022 If PRESERVE_DATE matches  \u2192 skip (mark as PRESERVED)",
        "  \u2022 Otherwise \u2192 call AGENT:",
        "    PKG_AGENT_ARCHIVE.PRC_CLEANUP_UNIT('TRUNCATE', ...)",
        "  \u2022 Set TRUNCATE_STATUS = 'Y'",
        "",
        "PREVIEW MODE: p_truncate_execute = 'N' -- show what would be done",
        "EXECUTE MODE: p_truncate_execute = 'Y' -- actual TRUNCATE"
    ],
    subtitle="Preview first -- approve before you execute",
    note="PKG_ARCHIVER_TRUNCATE.prc_truncate(p_execute, p_target_owner, p_target_table_name)"
)

# ── SLIDE 19: Pipeline Flow Diagram ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, WHITE)
add_title_header(slide, "ARCHIVER Pipeline Flow",
                 "DISCOVER \u2192 ARCHIVE \u2192 QUALITY \u2192 TRUNCATE preview")
add_pipeline_steps_diagram(slide, Inches(1.0), Inches(1.5), Inches(7.0))
add_pipeline_flow(slide, Inches(1.5), Inches(6.3), Inches(10), [
    ("DISCOVER", AGENT_C), ("ARCHIVE", ARCH_C), ("QUALITY", ORANGE), ("TRUNCATE", RED)
], arrow_color=NAVY)
_add_textbox(slide, Inches(0.6), Inches(6.0), Inches(5), Inches(0.3),
             "Flow overview", font_size=13, bold=True, color=NAVY)
add_bottom_bar(slide)
add_slide_number(slide, len(prs.slides))

# ── SLIDE 20: Runner Examples ──
slide_full_code(
    "Runner Examples",
    "-- DISCOVER only for a single table\n"
    "BEGIN\n"
    "  PKG_ARCHIVER_RUNNER.prc_run_table(\n"
    "    p_source_db_link   => 'AGENT_LINK',\n"
    "    p_owner            => 'CLIENT1',\n"
    "    p_table_name       => 'ORDERS_SUBPART_SRC',\n"
    "    p_execute          => 'Y',\n"
    "    p_stop_after_step  => 'DISCOVER',\n"
    "    p_truncate_execute => 'N'\n"
    "  );\n"
    "END;\n"
    "/\n"
    "\n"
    "-- Full pipeline for all enabled sources\n"
    "BEGIN\n"
    "  PKG_ARCHIVER_RUNNER.prc_run_all(\n"
    "    p_execute          => 'Y',\n"
    "    p_stop_after_step  => 'TRUNCATE',\n"
    "    p_truncate_execute => 'N'   -- preview only\n"
    "  );\n"
    "END;\n"
    "/\n"
    "\n"
    "-- With recommended TRUNCATE preview\n"
    "BEGIN\n"
    "  PKG_ARCHIVER_RUNNER.prc_run_table(\n"
    "    p_source_db_link   => 'AGENT_LINK',\n"
    "    p_owner            => 'CLIENT1',\n"
    "    p_table_name       => 'ORDERS_ARCH_SRC',\n"
    "    p_execute          => 'Y',\n"
    "    p_stop_after_step  => 'TRUNCATE',\n"
    "    p_truncate_execute => 'N'\n"
    "  );\n"
    "END;\n"
    "/",
    subtitle="PKG_ARCHIVER_RUNNER -- pipeline orchestration",
    note="p_stop_after_step = DISCOVER | ARCHIVE | QUALITY | TRUNCATE"
)

# ── SLIDE 21: ARCHIVER Outputs ──
slide_two_columns(
    "ARCHIVER Outputs & Process Views",
    [
        "Tables:",
        "  \u2022 TBL_ARCHIVER_TABLES",
        "  \u2022 TBL_ARCHIVER_PARTITIONS",
        "  \u2022 TBL_ARCHIVER_RUNS",
        "  \u2022 TBL_ARCHIVER_PROCESS_LOG",
        "",
        "Unit status flags:",
        "  \u2022 ARCHIVE_STATUS (Y/N)",
        "  \u2022 QUALITY_STATUS (Y/N)",
        "  \u2022 TRUNCATE_STATUS (Y/N)"
    ],
    [
        "Process views:",
        "  \u2022 VW_ARCHIVER_SOURCE_PARTITIONS",
        "     \u2192 source metadata via DB link",
        "  \u2022 VW_ARCHIVER_DISCOVERY_PARTITIONS",
        "     \u2192 new partitions to discover",
        "  \u2022 VW_ARCHIVER_IMPORT_PARTITIONS",
        "     \u2192 ready for archive import",
        "  \u2022 VW_ARCHIVER_QUALITY_PARTITIONS",
        "     \u2192 awaiting quality check",
        "  \u2022 VW_ARCHIVER_TRUNCATE_PARTITIONS",
        "     \u2192 cleanup candidates"
    ],
    subtitle="What to inspect after running the pipeline"
)

# ── SLIDE 22: Post-Run Checks ──
slide_two_columns(
    "Post-Run Inspection",
    [
        "Runs (TBL_ARCHIVER_RUNS):",
        "  \u2022 RUN_STATUS (INIT/RUNNING/SUCCESS/WARNING/ERROR)",
        "  \u2022 RUN_TYPE (DISCOVER/ARCHIVE/QUALITY/TRUNCATE/RUNNER)",
        "  \u2022 STARTED_AT / ENDED_AT",
        "  \u2022 ERROR_MESSAGE",
        "",
        "Logs (TBL_ARCHIVER_PROCESS_LOG):",
        "  \u2022 Detailed messages per step",
        "  \u2022 Summaries between markers:",
        "    <<<PARTMGR_SUMMARY_BEGIN/END>>>",
        "  \u2022 Generated SQL from preview mode"
    ],
    [
        "Partitions (TBL_ARCHIVER_PARTITIONS):",
        "  \u2022 ARCHIVE_STATUS -- loaded or not",
        "  \u2022 QUALITY_STATUS -- row count match",
        "  \u2022 TRUNCATE_STATUS -- source cleaned or not",
        "  \u2022 SOURCE_ROW_COUNT vs TARGET_ROW_COUNT",
        "  \u2022 ERROR_MESSAGE -- failure reasons",
        "",
        "Preview queries:",
        "  SELECT * FROM VW_ARCHIVER_IMPORT_PARTITIONS;",
        "  SELECT * FROM VW_ARCHIVER_QUALITY_PARTITIONS;",
        "  SELECT * FROM VW_ARCHIVER_TRUNCATE_PARTITIONS;"
    ],
    subtitle="What to verify after each pipeline step"
)

# ── SLIDE 23: Reports ──
slide_content(
    "ARCHIVER Reports",
    [
        "Send a report:",
        "  PKG_UTIL_MAIL.prc_send_report('ARCHIVER_SUMMARY');",
        "",
        "HTML report contents:",
        "  \u2022 Report window (REPORT_LOOKBACK_DAYS, default 7)",
        "  \u2022 Executive summary -- run counts and latest statuses",
        "  \u2022 Data status -- ARCHIVE/QUALITY/TRUNCATE progress",
        "  \u2022 Pending work -- what still needs attention",
        "  \u2022 Latest warnings and errors",
        "  \u2022 Full process summaries as HTML attachments",
        "",
        "Configuration (TBL_UTIL_CONFIG):",
        "  REPORT_LOOKBACK_DAYS, REPORT_SUMMARY_MAX_CHARS, REPORT_MAX_ROWS"
    ],
    subtitle="HTML email reports with attachments",
    note="PKG_UTIL_MAIL uses UTL_SMTP for plain SMTP delivery."
)

# ── SLIDE 24: Safety ──
slide_content(
    "Safety Rules",
    [
        "SOURCE_DB_LINK must be real and non-null -- LOCAL and NONE are blocked",
        "All object names are validated via DBMS_ASSERT before dynamic SQL",
        "Cleanup (TRUNCATE/PURGE) defaults to preview mode",
        "TRUNCATE requires explicit p_truncate_execute = 'Y'",
        "PURGE requires p_purge_execute = 'Y'",
        "Quality check must pass before cleanup candidates are selected",
        "All runs and process logs are fully auditable",
        "",
        "Golden rule:",
        "  discover \u2192 archive \u2192 quality \u2192 preview truncate \u2192 approve \u2192 execute"
    ],
    subtitle="Safety first, always"
)

# ── SLIDE 25: REPLICA Flow (with diagram) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, WHITE)
add_title_header(slide, "REPLICA -- Data Flow",
                 "DISCOVER \u2192 REPLICATE \u2192 QUALITY \u2192 PURGE preview")
add_replica_flow_diagram(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.0))
add_pipeline_flow(slide, Inches(1.0), Inches(5.5), Inches(11), [
    ("DISCOVER", AGENT_C), ("REPLICATE", REPL_C), ("QUALITY", ORANGE), ("PURGE", RED)
], arrow_color=NAVY)
add_bullets(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(1.5), [
    "REPLICA reads qualified ARCHIVER metadata via REPLICA_ARCHIVER_PARTITIONS_SRC (synonym through DB link)",
    "Uses staging + EXCHANGE PARTITION (same pattern as ARCHIVER)",
    "PURGE operates only on local REPLICA targets -- does not affect ARCHIVER or CLIENT",
], font_size=14)
add_bottom_bar(slide)
add_slide_number(slide, len(prs.slides))

# ── SLIDE 26: REPLICA Purpose ──
slide_content(
    "REPLICA -- Purpose & Use Cases",
    [
        "REPLICA is an optional Layer 3 component",
        "Use it when archived data must be copied further downstream:",
        "  \u2022 Reporting database",
        "  \u2022 Bounded operational copy",
        "  \u2022 Isolated archive consumer",
        "",
        "REPLICA reads ARCHIVER through a real DB link",
        "It does not need direct access to CLIENT schemas or AGENT",
        "",
        "Base system: CLIENT \u2192 AGENT \u2192 ARCHIVER",
        "Extended system: + ARCHIVER \u2192 REPLICA"
    ],
    subtitle="Optional downstream layer"
)

# ── SLIDE 27: REPLICA Config Example ──
slide_with_code(
    "REPLICA -- Configuration Example",
    [
        "Source for REPLICA:",
        "  PARTMGR_ARCHIVER.TBL_ARCHIVER_CLIENT1_SUBPART",
        "  via ARCHIVER_LINK",
        "",
        "Target: PARTMGR_REPLICA.TBL_REPLICA_CLIENT1_SUBPART",
        "",
        "REPLICA depends on:",
        "  \u2022 ARCHIVER metadata (TBL_ARCHIVER_PARTITIONS)",
        "  \u2022 Physical data in ARCHIVER target tables",
        "",
        "DAYS_ONLINE = 365",
        "  \u2192 replicated data is retained for one year"
    ],
    "-- TBL_REPLICA_TABLES row\n"
    "INSERT INTO TBL_REPLICA_TABLES (\n"
    "  SOURCE_DB_LINK,\n"
    "  SOURCE_OWNER,\n"
    "  SOURCE_TABLE_NAME,\n"
    "  TARGET_OWNER,\n"
    "  TARGET_TABLE_NAME,\n"
    "  PARALLEL_DEGREE,\n"
    "  DAYS_ONLINE\n"
    ") VALUES (\n"
    "  'ARCHIVER_LINK',\n"
    "  'PARTMGR_ARCHIVER',\n"
    "  'TBL_ARCHIVER_CLIENT1_SUBPART',\n"
    "  'PARTMGR_REPLICA',\n"
    "  'TBL_REPLICA_CLIENT1_SUBPART',\n"
    "  4,\n"
    "  365\n"
    ");",
    note="REPLICA_ARCHIVER_PARTITIONS_SRC is a synonym pointing to TBL_ARCHIVER_PARTITIONS@ARCHIVER_LINK"
)

# ── SLIDE 28: REPLICA Execution ──
slide_full_code(
    "REPLICA -- Execution",
    "-- Full REPLICA pipeline with PURGE preview\n"
    "BEGIN\n"
    "  PKG_REPLICA_RUNNER.prc_run(\n"
    "    p_execute         => 'Y',\n"
    "    p_stop_after_step => 'PURGE',\n"
    "    p_purge_execute   => 'N'    -- preview only\n"
    "  );\n"
    "END;\n"
    "/\n"
    "\n"
    "-- Stop after REPLICATE (skip quality and purge)\n"
    "BEGIN\n"
    "  PKG_REPLICA_RUNNER.prc_run(\n"
    "    p_execute         => 'Y',\n"
    "    p_stop_after_step => 'REPLICATE',\n"
    "    p_purge_execute   => 'N'\n"
    "  );\n"
    "END;\n"
    "/",
    subtitle="PKG_REPLICA_RUNNER.prc_run -- orchestration",
    note="Default pattern: replicate + quality execute, purge preview."
)

# ── SLIDE 29: REPLICA Outputs ──
slide_two_columns(
    "REPLICA -- Outputs & Reports",
    [
        "Tables:",
        "  \u2022 TBL_REPLICA_TABLES",
        "  \u2022 TBL_REPLICA_PARTITIONS",
        "  \u2022 TBL_REPLICA_RUNS",
        "  \u2022 TBL_REPLICA_PROCESS_LOG",
        "",
        "Report:",
        "  PKG_UTIL_MAIL.prc_send_report(",
        "    'REPLICA_SUMMARY');"
    ],
    [
        "Process views:",
        "  \u2022 VW_REPLICA_DISCOVERY_PARTITIONS",
        "  \u2022 VW_REPLICA_REPLICATE_PARTITIONS",
        "  \u2022 VW_REPLICA_QUALITY_PARTITIONS",
        "  \u2022 VW_REPLICA_PURGE_PARTITIONS",
        "",
        "Safety:",
        "  PURGE requires p_purge_execute = 'Y'"
    ],
    subtitle="Same patterns as ARCHIVER, own data set"
)

# ── SLIDE 30: Operating Pattern ──
slide_content(
    "Operating Pattern -- Step by Step",
    [
        "1. Configure the source table in TBL_ARCHIVER_TABLES",
        "2. Run DISCOVER -- detect new partitions",
        "3. Review TBL_ARCHIVER_PARTITIONS",
        "4. Run ARCHIVE + QUALITY",
        "5. Verify row counts and process logs",
        "6. Preview truncate candidates (VW_ARCHIVER_TRUNCATE_PARTITIONS)",
        "7. Approve and execute TRUNCATE with p_truncate_execute = 'Y'",
        "8. Optionally run REPLICA",
        "9. Send summary reports",
        "",
        "Operational golden rule:",
        "  configure \u2192 discover \u2192 archive \u2192 quality \u2192 preview \u2192 approve \u2192 report"
    ],
    subtitle="A typical operational cycle"
)

# ── SLIDE 31: SPLIT Topology (diagram) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, WHITE)
add_title_header(slide, "Installation Topology -- SPLIT Model",
                 "Three separate databases, real DB links between components")
add_splited_topology_diagram(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(3.5))
add_bullets(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5), [
    "AGENT installed on the SOURCE database (alongside CLIENT schemas)",
    "ARCHIVER installed on a dedicated ARCHIVE database",
    "REPLICA installed on a dedicated REPLICA database",
    "Components communicate through real Oracle database links",
    "Three separate SYS connections manage schema lifecycle (SOURCE_SYS_CONNECT, ARCHIVER_SYS_CONNECT, REPLICA_SYS_CONNECT)",
], font_size=14)
add_bottom_bar(slide)
add_slide_number(slide, len(prs.slides))

# ── SLIDE 32: Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, NAVY)
_add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Summary", font_size=40, bold=True, color=ORANGE,
             alignment=PP_ALIGN.CENTER)

summary_items = [
    "PARTMGR is a centralized Oracle partition archive control plane in a 3-layer architecture",
    "",
    "AGENT -- thin helper layer at the source (policy-free)",
    "ARCHIVER -- main control: DISCOVER \u2192 ARCHIVE \u2192 QUALITY \u2192 TRUNCATE preview",
    "REPLICA -- optional downstream replica: DISCOVER \u2192 REPLICATE \u2192 QUALITY \u2192 PURGE preview",
    "",
    "Key features:",
    "  \u2022 HIGH_VALUE-based identity -- independent of physical names",
    "  \u2022 Preview-first -- destructive operations require explicit approval",
    "  \u2022 Quality gate -- no cleanup without row-count match",
    "  \u2022 Real DB links -- no magic LOCAL/NONE values",
    "  \u2022 Built-in HTML email reports",
    "",
    "Further reading: docs/architecture.md \u00b7 docs/installation.md \u00b7 docs/operations.md"
]
add_bullets(slide, Inches(1), Inches(1.8), Inches(11), Inches(4.8),
            summary_items, font_size=16, color=WHITE)
add_bottom_bar(slide)
add_slide_number(slide, 31)

# ── Save ──
output_path = "PARTMGR_Prezentacja.pptx"
prs.save(output_path)
print(f"Done -- saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
